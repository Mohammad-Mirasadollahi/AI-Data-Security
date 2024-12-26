# document_loader.py

import configparser
import csv
import json
import logging
import os
import xml.etree.ElementTree as ET
from typing import Tuple, List, Dict

import textract
import yaml
from PyPDF2 import PdfReader
from bs4 import BeautifulSoup
from docx import Document
from ebooklib import epub
from pptx import Presentation

from Utils.utils import calculate_sha256, calculate_fuzzy_hash


class DocumentLoader:
    """
    A class to load and extract text from various types of text-based files,
    and compute their SHA-256 and Fuzzy hashes.
    """

    # Supported extensions mapped to their extraction methods
    SUPPORTED_EXTENSIONS = {
        '.txt': 'extract_txt',
        '.pdf': 'extract_pdf',
        '.docx': 'extract_docx',
        '.rtf': 'extract_rtf',
        '.pptx': 'extract_pptx',
        '.html': 'extract_html',
        '.htm': 'extract_html',
        '.md': 'extract_md',
        '.csv': 'extract_csv',
        '.epub': 'extract_epub',
        '.json': 'extract_json',
        '.xml': 'extract_xml',
        '.yaml': 'extract_yaml',
        '.yml': 'extract_yaml',
        '.ini': 'extract_ini',
        '.log': 'extract_log',
        '.sql': 'extract_sql',
    }

    def __init__(self, log_file: str = 'document_loader.log'):
        """
        Initializes the DocumentLoader with logging configurations.

        Parameters:
        log_file (str): The name of the log file.
        """
        self.file_names = []
        self.documents = []
        self.logger = logging.getLogger('DocumentLoader')
        self.logger.setLevel(logging.INFO)
        handler = logging.FileHandler(log_file)
        handler.setFormatter(logging.Formatter('%(asctime)s %(levelname)s:%(message)s'))
        if not self.logger.handlers:
            self.logger.addHandler(handler)
        self.logger.info("DocumentLoader initialized.")

    # Extraction methods for each supported file type

    def extract_txt(self, file_path):
        """Extracts text from a .txt file."""
        with open(file_path, 'r', encoding='utf-8', errors='replace') as f:
            return f.read()

    def extract_pdf(self, file_path):
        """Extracts text from a .pdf file."""
        reader = PdfReader(file_path)
        text = []
        for page in reader.pages:
            page_text = page.extract_text()
            if page_text:
                text.append(page_text)
        return '\n'.join(text)

    def extract_docx(self, file_path):
        """Extracts text from a .docx file."""
        doc = Document(file_path)
        return '\n'.join([para.text for para in doc.paragraphs])

    def extract_rtf(self, file_path):
        """Extracts text from a .rtf file."""
        try:
            text = textract.process(file_path, encoding='utf-8', errors='replace').decode('utf-8', errors='replace')
            return text
        except Exception as e:
            self.logger.error(f"Failed to extract RTF from '{file_path}': {e}")
            return ""

    def extract_pptx(self, file_path):
        """Extracts text from a .pptx file."""
        prs = Presentation(file_path)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, 'text'):
                    text.append(shape.text)
        return '\n'.join(text)

    def extract_html(self, file_path):
        """Extracts text from a .html or .htm file."""
        try:
            with open(file_path, 'r', encoding='utf-8', errors='replace') as f:
                soup = BeautifulSoup(f, 'html.parser')
            return soup.get_text()
        except Exception as e:
            self.logger.error(f"Failed to extract HTML from '{file_path}': {e}")
            return ""

    def extract_md(self, file_path):
        """Extracts text from a .md (Markdown) file."""
        with open(file_path, 'r', encoding='utf-8', errors='replace') as f:
            return f.read()

    def extract_csv(self, file_path):
        """Extracts text from a .csv file."""
        try:
            with open(file_path, 'r', encoding='utf-8', errors='replace') as f:
                reader = csv.reader(f)
                return '\n'.join(['\t'.join(row) for row in reader])
        except Exception as e:
            self.logger.error(f"Failed to extract CSV from '{file_path}': {e}")
            return ""

    def extract_epub(self, file_path):
        """Extracts text from a .epub file."""
        try:
            book = epub.read_epub(file_path)
            text = []
            for item in book.get_items():
                if item.get_type() == epub.ITEM_DOCUMENT:
                    soup = BeautifulSoup(item.get_content(), 'html.parser')
                    text.append(soup.get_text())
            return '\n'.join(text)
        except Exception as e:
            self.logger.error(f"Failed to extract EPUB from '{file_path}': {e}")
            return ""

    def extract_json(self, file_path):
        """Extracts text from a .json file."""
        try:
            with open(file_path, 'r', encoding='utf-8', errors='replace') as f:
                data = json.load(f)
            return json.dumps(data, ensure_ascii=False, indent=4)
        except Exception as e:
            self.logger.error(f"Failed to extract JSON from '{file_path}': {e}")
            return ""

    def extract_xml(self, file_path):
        """Extracts text from a .xml file."""
        try:
            tree = ET.parse(file_path)
            root = tree.getroot()
            return self._extract_text_from_xml(root)
        except Exception as e:
            self.logger.error(f"Failed to extract XML from '{file_path}': {e}")
            return ""

    def _extract_text_from_xml(self, element):
        """Recursively extracts text from XML elements."""
        texts = []
        if element.text:
            texts.append(element.text)
        for child in element:
            texts.append(self._extract_text_from_xml(child))
            if child.tail:
                texts.append(child.tail)
        return ''.join(texts)

    def extract_yaml(self, file_path):
        """Extracts text from a .yaml or .yml file."""
        try:
            with open(file_path, 'r', encoding='utf-8', errors='replace') as f:
                data = yaml.safe_load(f)
            return yaml.dump(data, allow_unicode=True, sort_keys=False)
        except Exception as e:
            self.logger.error(f"Failed to extract YAML from '{file_path}': {e}")
            return ""

    def extract_ini(self, file_path):
        """Extracts text from a .ini file."""
        try:
            config = configparser.ConfigParser()
            config.read(file_path, encoding='utf-8')
            output = []
            for section in config.sections():
                output.append(f"[{section}]")
                for key, value in config.items(section):
                    output.append(f"{key} = {value}")
                output.append("")  # Add a newline between sections
            return '\n'.join(output)
        except Exception as e:
            self.logger.error(f"Failed to extract INI from '{file_path}': {e}")
            return ""

    def extract_log(self, file_path):
        """Extracts text from a .log file."""
        try:
            with open(file_path, 'r', encoding='utf-8', errors='replace') as f:
                return f.read()
        except Exception as e:
            self.logger.error(f"Failed to extract LOG from '{file_path}': {e}")
            return ""

    def extract_sql(self, file_path):
        """Extracts text from a .sql file."""
        try:
            with open(file_path, 'r', encoding='utf-8', errors='replace') as f:
                return f.read()
        except Exception as e:
            self.logger.error(f"Failed to extract SQL from '{file_path}': {e}")
            return ""

    # Methods to add or remove supported extensions dynamically

    def add_supported_extension(self, extension, method_name):
        """
        Adds a new supported file extension with its extraction method.

        Parameters:
        extension (str): File extension (e.g., '.json').
        method_name (str): Name of the extraction method (e.g., 'extract_json').

        Raises:
        AttributeError: If the extraction method does not exist in the class.
        """
        if hasattr(self, method_name):
            self.SUPPORTED_EXTENSIONS[extension.lower()] = method_name
            self.logger.info(f"Added support for {extension} with method {method_name}.")
        else:
            self.logger.error(f"Method {method_name} does not exist in DocumentLoader.")

    def remove_supported_extension(self, extension):
        """
        Removes a supported file extension.

        Parameters:
        extension (str): File extension to remove (e.g., '.json').
        """
        if extension.lower() in self.SUPPORTED_EXTENSIONS:
            del self.SUPPORTED_EXTENSIONS[extension.lower()]
            self.logger.info(f"Removed support for {extension}.")
        else:
            self.logger.warning(f"Extension {extension} is not supported and cannot be removed.")

    def load_documents(self, folder_path: str) -> Tuple[List[str], List[Dict]]:
        """
        Loads and extracts text from all supported documents in a specified folder,
        and computes their SHA-256 and Fuzzy hashes.

        Parameters:
        folder_path (str): Path to the folder containing documents.

        Returns:
        Tuple[List[str], List[Dict]]: A tuple containing a list of file names and a list of their corresponding extracted texts with hashes.
        """
        self.file_names = []
        self.documents = []

        if not os.path.isdir(folder_path):
            self.logger.error(f"Invalid folder path: {folder_path}")
            return self.file_names, self.documents

        for filename in os.listdir(folder_path):
            file_path = os.path.join(folder_path, filename)
            file_ext = os.path.splitext(filename)[1].lower()

            if os.path.isfile(file_path) and file_ext in self.SUPPORTED_EXTENSIONS:
                try:
                    method_name = self.SUPPORTED_EXTENSIONS[file_ext]
                    extractor = getattr(self, method_name)
                    text = extractor(file_path)
                    if text.strip():  # Ensure that extracted text is not empty
                        sha256 = calculate_sha256(file_path)
                        fuzzy = calculate_fuzzy_hash(file_path)
                        self.documents.append({
                            "file_name": filename,
                            "file_type": file_ext.strip('.'),
                            "text": text,
                            "sha256": sha256,
                            "fuzzy_hash": fuzzy
                        })
                        self.file_names.append(filename)
                        self.logger.info(f"Successfully loaded file: {filename}")
                    else:
                        self.logger.warning(f"No text extracted from file: {filename}")
                except Exception as e:
                    self.logger.error(f"Error reading file {filename}: {e}")
            else:
                self.logger.info(f"Unsupported or non-file: {filename}")

        self.logger.info(f"Total documents loaded: {len(self.documents)}")
        return self.file_names, self.documents

    def get_documents(self) -> Tuple[List[str], List[Dict]]:
        """
        Retrieves the loaded documents.

        Returns:
        Tuple[List[str], List[Dict]]: A tuple containing a list of file names and a list of their corresponding extracted texts with hashes.
        """
        return self.file_names, self.documents
